'''PowerPoint ファイル (.pptx) を走査'''

from pathlib import Path
from typing import List, Tuple, cast

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore
from pptx.text.text import TextFrame
from typing import cast

from .base import BaseScanner, Finding
from .unicode_utils import (CONTRAST_THRESHOLD, MIN_FONT_PT,
                            contains_invisible_unicode, contrast_ratio)

HIDDEN_ATTR = "{http://schemas.openxmlformats.org/presentationml/2006/main}show"



class PptxScanner(BaseScanner):
    '''.pptx 専用スキャナ'''

    def scan(self, path: Path) -> Tuple[bool, List[Finding]]:
        prs = Presentation(str(path))
        findings: List[Finding] = []

        for s_idx, slide in enumerate(prs.slides, 1):
            # 非表示スライドチェック slide/@show="0" :contentReference[oaicite:4]{index=4}
            if slide._element.get(HIDDEN_ATTR, None) == "0":
                findings.append(Finding(
                    location=f"スライド {s_idx} (非表示)",
                    snippet="このスライドはスライドショーで非表示設定"
                ))

            # 各 shape のテキスト走査
            for shape in slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    continue
                text_frame = cast(TextFrame, shape.text_frame)  # type: ignore[attr-defined]
                
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        txt = run.text or ""
                        if not txt.strip():
                            continue

                        # 極小フォント
                        size_pt = run.font.size.pt if run.font.size else 18
                        tiny = size_pt <= MIN_FONT_PT

                        # 文字色 (python-pptx は alpha 未対応) :contentReference[oaicite:5]{index=5}
                        if run.font.color.type is None:
                            rgb = (0, 0, 0)
                        else:
                            rgb = tuple(run.font.color.rgb[i] for i in range(3))  # type: ignore

                        low_contrast = contrast_ratio(rgb) < CONTRAST_THRESHOLD

                        invisible_unicode = contains_invisible_unicode(txt)

                        if any([tiny, low_contrast, invisible_unicode]):
                            snippet = (txt[:40]
                                       + (f" …+{len(txt)-40}文字" if len(txt) > 40 else ""))
                            findings.append(Finding(
                                location=f"スライド {s_idx}",
                                snippet=snippet
                            ))

        return (bool(findings), findings)
